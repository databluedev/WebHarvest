"""Document extraction service for PDF, DOCX, XLSX, PPTX, CSV, RTF, EPUB and other non-HTML formats."""

import base64
import csv
import io
import logging
import zipfile
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)


@dataclass
class DocumentResult:
    """Extracted document content."""
    text: str = ""
    markdown: str = ""
    metadata: dict = field(default_factory=dict)
    page_count: int = 0
    word_count: int = 0
    tables: list[dict] = field(default_factory=list)    # [{headers, rows, page}]
    images: list[dict] = field(default_factory=list)     # [{data_b64, page, width, height}]
    pages: list[dict] = field(default_factory=list)      # [{page_num, text, markdown}]


def detect_document_type(
    url: str,
    content_type: str | None = None,
    raw_bytes: bytes = b"",
) -> str:
    """Detect document type from URL extension, content-type header, and magic bytes.

    Returns: "html", "pdf", "docx", "xlsx", "pptx", "csv", "rtf", "epub", or "unknown"
    """
    url_lower = url.lower().split("?")[0].split("#")[0]

    # Check URL extension first
    ext_map = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "docx",
        ".xlsx": "xlsx",
        ".xls": "xlsx",
        ".pptx": "pptx",
        ".ppt": "pptx",
        ".csv": "csv",
        ".tsv": "csv",
        ".rtf": "rtf",
        ".epub": "epub",
    }
    for ext, doc_type in ext_map.items():
        if url_lower.endswith(ext):
            return doc_type

    # Check content-type header
    if content_type:
        ct = content_type.lower()
        ct_map = {
            "application/pdf": "pdf",
            "wordprocessingml": "docx",
            "application/msword": "docx",
            "spreadsheetml": "xlsx",
            "application/vnd.ms-excel": "xlsx",
            "presentationml": "pptx",
            "application/vnd.ms-powerpoint": "pptx",
            "text/csv": "csv",
            "application/rtf": "rtf",
            "text/rtf": "rtf",
            "application/epub": "epub",
        }
        for pattern, doc_type in ct_map.items():
            if pattern in ct:
                return doc_type
        if "text/html" in ct or "application/xhtml" in ct:
            return "html"

    # Check magic bytes
    if raw_bytes:
        if raw_bytes[:4] == b"%PDF":
            return "pdf"
        if raw_bytes[:5] == b"{\\rtf":
            return "rtf"
        if raw_bytes[:4] == b"\xd0\xcf\x11\xe0":
            # OLE2 compound document — legacy xls/doc/ppt
            return "xlsx"  # Most common OLE2 in web context
        if raw_bytes[:4] == b"PK\x03\x04":
            # ZIP-based format — inspect contents to distinguish
            return _detect_zip_format(raw_bytes)

    return "html"  # Default to HTML


def _detect_zip_format(raw_bytes: bytes) -> str:
    """Inspect ZIP archive contents to distinguish docx/xlsx/pptx/epub."""
    try:
        with zipfile.ZipFile(io.BytesIO(raw_bytes)) as zf:
            names = zf.namelist()
            if any(n.startswith("word/") for n in names):
                return "docx"
            if any(n.startswith("xl/") for n in names):
                return "xlsx"
            if any(n.startswith("ppt/") for n in names):
                return "pptx"
            if "META-INF/container.xml" in names:
                return "epub"
    except Exception:
        pass
    return "docx"  # Default ZIP to docx


# ---------------------------------------------------------------------------
# PDF Extraction — Multi-Strategy
# ---------------------------------------------------------------------------

def _table_to_markdown(headers: list[str], rows: list[list[str]]) -> str:
    """Convert table data to a markdown pipe table."""
    if not headers and not rows:
        return ""
    if not headers and rows:
        headers = [f"Col {i+1}" for i in range(len(rows[0]))]
    lines = []
    lines.append("| " + " | ".join(str(h) for h in headers) + " |")
    lines.append("| " + " | ".join("---" for _ in headers) + " |")
    for row in rows:
        # Pad or trim row to match header count
        padded = list(row) + [""] * (len(headers) - len(row))
        lines.append("| " + " | ".join(str(c) for c in padded[:len(headers)]) + " |")
    return "\n".join(lines)


def _extract_pdf_tables_pymupdf(page) -> list[dict]:
    """Extract tables from a PyMuPDF page using find_tables() (PyMuPDF 1.23+)."""
    tables = []
    try:
        tab_finder = page.find_tables()
        for tab in tab_finder.tables:
            data = tab.extract()
            if not data:
                continue
            headers = [str(c) if c else "" for c in data[0]]
            rows = [[str(c) if c else "" for c in row] for row in data[1:]]
            tables.append({"headers": headers, "rows": rows, "page": page.number + 1})
    except Exception:
        pass
    return tables


def _extract_pdf_tables_pdfplumber(raw_bytes: bytes, page_num: int) -> list[dict]:
    """Extract tables from a specific page using pdfplumber (fallback for complex layouts)."""
    try:
        import pdfplumber
    except ImportError:
        return []
    tables = []
    try:
        with pdfplumber.open(io.BytesIO(raw_bytes)) as pdf:
            if page_num < len(pdf.pages):
                plumber_page = pdf.pages[page_num]
                for tab in plumber_page.extract_tables():
                    if not tab:
                        continue
                    headers = [str(c) if c else "" for c in tab[0]]
                    rows = [[str(c) if c else "" for c in row] for row in tab[1:]]
                    tables.append({"headers": headers, "rows": rows, "page": page_num + 1})
    except Exception:
        pass
    return tables


def _extract_pdf_images(doc, page, max_images: int = 20, max_size: int = 2 * 1024 * 1024) -> list[dict]:
    """Extract embedded images from a PyMuPDF page."""
    images = []
    try:
        image_list = page.get_images(full=True)
        for img_info in image_list[:max_images - len(images)]:
            if len(images) >= max_images:
                break
            xref = img_info[0]
            try:
                img_data = doc.extract_image(xref)
                if not img_data or not img_data.get("image"):
                    continue
                raw = img_data["image"]
                if len(raw) > max_size:
                    continue
                images.append({
                    "data_b64": base64.b64encode(raw).decode("ascii"),
                    "page": page.number + 1,
                    "width": img_data.get("width", 0),
                    "height": img_data.get("height", 0),
                    "ext": img_data.get("ext", "png"),
                })
            except Exception:
                continue
    except Exception:
        pass
    return images


def _build_page_markdown(page, page_text: str, page_tables: list[dict], page_images: list[dict], image_offset: int) -> str:
    """Build enhanced markdown for a single PDF page using structured text analysis."""
    md_parts = []

    # Try structured text for heading detection
    try:
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if block.get("type") != 0:  # text block
                continue
            for line in block.get("lines", []):
                line_text = ""
                max_size = 0
                is_bold = False
                is_italic = False
                for span in line.get("spans", []):
                    span_text = span.get("text", "").strip()
                    if not span_text:
                        continue
                    line_text += span_text + " "
                    size = span.get("size", 12)
                    if size > max_size:
                        max_size = size
                    flags = span.get("flags", 0)
                    if flags & 2**4:  # bold
                        is_bold = True
                    if flags & 2**1:  # italic
                        is_italic = True

                line_text = line_text.strip()
                if not line_text:
                    continue

                if max_size >= 20:
                    md_parts.append(f"# {line_text}")
                elif max_size >= 16:
                    md_parts.append(f"## {line_text}")
                elif max_size >= 14:
                    md_parts.append(f"### {line_text}")
                elif is_bold and not is_italic:
                    md_parts.append(f"**{line_text}**")
                elif is_italic and not is_bold:
                    md_parts.append(f"*{line_text}*")
                elif is_bold and is_italic:
                    md_parts.append(f"***{line_text}***")
                else:
                    md_parts.append(line_text)
    except Exception:
        # Fallback to plain text if structured extraction fails
        md_parts.append(page_text.strip())

    # Add tables as markdown
    for table in page_tables:
        md_parts.append("")
        md_parts.append(_table_to_markdown(table["headers"], table["rows"]))
        md_parts.append("")

    # Add image placeholders
    for i, img in enumerate(page_images):
        md_parts.append(f"![Image {image_offset + i + 1} from page {img['page']}](embedded)")

    # Add links
    try:
        links = page.get_links()
        link_items = []
        for link in links:
            uri = link.get("uri", "")
            if uri:
                link_items.append(f"- [{uri}]({uri})")
        if link_items:
            md_parts.append("")
            md_parts.extend(link_items)
    except Exception:
        pass

    return "\n".join(md_parts)


async def extract_pdf(raw_bytes: bytes) -> DocumentResult:
    """Extract text, tables, images and metadata from a PDF using multi-strategy approach.

    Strategy 1: PyMuPDF (primary) — text, tables, images, structured text
    Strategy 2: pdfplumber (fallback for complex table layouts)
    Strategy 3: OCR fallback for scanned/image PDFs
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("PyMuPDF not installed, returning empty result")
        return DocumentResult(
            text="[PDF extraction requires PyMuPDF]",
            markdown="*PDF extraction requires PyMuPDF*",
            metadata={"error": "PyMuPDF not installed"},
        )

    try:
        doc = fitz.open(stream=raw_bytes, filetype="pdf")

        # Extract metadata
        meta = doc.metadata or {}
        metadata = {
            "author": meta.get("author", ""),
            "title": meta.get("title", ""),
            "subject": meta.get("subject", ""),
            "creator": meta.get("creator", ""),
            "producer": meta.get("producer", ""),
            "creation_date": meta.get("creationDate", ""),
            "mod_date": meta.get("modDate", ""),
            "page_count": doc.page_count,
            "document_type": "pdf",
        }

        # Extract table of contents
        toc = doc.get_toc()
        if toc:
            metadata["table_of_contents"] = [
                {"level": level, "title": title, "page": page}
                for level, title, page in toc
            ]

        all_tables = []
        all_images = []
        all_pages = []
        pages_text = []
        md_parts = []
        total_chars = 0

        # Build header
        title = metadata.get("title") or "PDF Document"
        md_parts.append(f"# {title}\n")
        if metadata.get("author"):
            md_parts.append(f"**Author:** {metadata['author']}\n")

        # Strategy 1: PyMuPDF extraction per page
        for page_num in range(doc.page_count):
            page = doc[page_num]
            text = page.get_text("text")
            total_chars += len(text.strip())
            pages_text.append(text)

            # Extract tables (Strategy 1: PyMuPDF)
            page_tables = _extract_pdf_tables_pymupdf(page)

            # Strategy 2: pdfplumber fallback if PyMuPDF found no tables
            if not page_tables:
                page_tables = _extract_pdf_tables_pdfplumber(raw_bytes, page_num)

            all_tables.extend(page_tables)

            # Extract images (limit total to 20)
            if len(all_images) < 20:
                page_images = _extract_pdf_images(doc, page, max_images=20 - len(all_images))
                all_images.extend(page_images)
            else:
                page_images = []

            # Build per-page markdown
            page_md = _build_page_markdown(page, text, page_tables, page_images, len(all_images) - len(page_images))

            all_pages.append({
                "page_num": page_num + 1,
                "text": text.strip(),
                "markdown": page_md,
            })

        # Strategy 3: OCR fallback for scanned PDFs
        avg_chars = total_chars / max(doc.page_count, 1)
        if avg_chars < 50 and doc.page_count > 0:
            metadata["scanned_pdf"] = True
            ocr_pages_text = []
            for page_num in range(doc.page_count):
                page = doc[page_num]
                try:
                    tp = page.get_textpage_ocr(language="eng")
                    ocr_text = page.get_text("text", textpage=tp)
                    if ocr_text.strip():
                        ocr_pages_text.append(ocr_text)
                        all_pages[page_num]["text"] = ocr_text.strip()
                        all_pages[page_num]["markdown"] = ocr_text.strip()
                except Exception:
                    # OCR unavailable — try pixmap render note
                    try:
                        _pix = page.get_pixmap(dpi=150)
                        metadata["ocr_unavailable"] = True
                        metadata["note"] = "Scanned PDF detected but OCR (Tesseract) is not available"
                    except Exception:
                        pass
            if ocr_pages_text:
                pages_text = ocr_pages_text
                metadata["ocr_applied"] = True

        full_text = "\n\n".join(t for t in pages_text if t.strip())
        word_count = len(full_text.split())

        # Finalize markdown
        md_parts.append(f"**Pages:** {doc.page_count} | **Words:** {word_count}\n")
        md_parts.append("---\n")
        for pg in all_pages:
            md_parts.append(f"## Page {pg['page_num']}\n")
            md_parts.append(pg["markdown"])
            md_parts.append("\n---\n")

        markdown = "\n\n".join(md_parts)

        metadata["table_count"] = len(all_tables)
        metadata["image_count"] = len(all_images)

        # Capture page_count BEFORE closing to avoid "document closed" bug
        total_pages = doc.page_count
        doc.close()

        return DocumentResult(
            text=full_text,
            markdown=markdown,
            metadata=metadata,
            page_count=total_pages,
            word_count=word_count,
            tables=all_tables,
            images=all_images,
            pages=all_pages,
        )

    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return DocumentResult(
            text=f"[PDF extraction failed: {e}]",
            markdown=f"*PDF extraction failed: {e}*",
            metadata={"error": str(e), "document_type": "pdf"},
        )


# ---------------------------------------------------------------------------
# DOCX Extraction — Enhanced
# ---------------------------------------------------------------------------

async def extract_docx(raw_bytes: bytes) -> DocumentResult:
    """Extract text, images, hyperlinks, and metadata from a DOCX."""
    try:
        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
    except ImportError:
        logger.warning("python-docx not installed, returning empty result")
        return DocumentResult(
            text="[DOCX extraction requires python-docx]",
            markdown="*DOCX extraction requires python-docx*",
            metadata={"error": "python-docx not installed"},
        )

    try:
        doc = Document(io.BytesIO(raw_bytes))

        # Extract metadata
        props = doc.core_properties
        metadata = {
            "author": props.author or "",
            "title": props.title or "",
            "subject": props.subject or "",
            "created": props.created.isoformat() if props.created else "",
            "modified": props.modified.isoformat() if props.modified else "",
            "last_modified_by": props.last_modified_by or "",
            "revision": props.revision,
            "document_type": "docx",
        }

        md_parts = []
        text_parts = []
        all_images = []

        title = metadata.get("title") or "Document"
        md_parts.append(f"# {title}\n")
        if metadata.get("author"):
            md_parts.append(f"**Author:** {metadata['author']}\n")
        md_parts.append("---\n")

        # Build hyperlink map from relationships
        hyperlink_map = {}
        try:
            for rel in doc.part.rels.values():
                if "hyperlink" in rel.reltype:
                    hyperlink_map[rel.rId] = rel.target_ref
        except Exception:
            pass

        # Extract paragraphs with enhanced formatting
        list_level_counter = {}
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                md_parts.append("")
                continue

            text_parts.append(text)

            style_name = para.style.name if para.style else ""

            # Detect list indentation level
            indent_level = 0
            try:
                pPr = para._element.find(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}ilvl")
                if pPr is not None:
                    indent_level = int(pPr.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", "0"))
            except Exception:
                pass

            # Build inline-formatted text from runs
            formatted_text = _build_formatted_runs(para, hyperlink_map)

            if style_name.startswith("Heading 1"):
                md_parts.append(f"# {formatted_text}")
            elif style_name.startswith("Heading 2"):
                md_parts.append(f"## {formatted_text}")
            elif style_name.startswith("Heading 3"):
                md_parts.append(f"### {formatted_text}")
            elif style_name.startswith("Heading 4"):
                md_parts.append(f"#### {formatted_text}")
            elif "List Bullet" in style_name or "List" in style_name and "Number" not in style_name:
                indent = "  " * indent_level
                md_parts.append(f"{indent}- {formatted_text}")
            elif "List Number" in style_name:
                indent = "  " * indent_level
                md_parts.append(f"{indent}1. {formatted_text}")
            else:
                md_parts.append(formatted_text)

        # Extract tables
        all_tables = []
        for table in doc.tables:
            md_parts.append("")
            table_headers = []
            table_rows = []
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                md_parts.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    md_parts.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    table_headers = cells
                else:
                    table_rows.append(cells)
                for cell_text in cells:
                    if cell_text:
                        text_parts.append(cell_text)
            all_tables.append({"headers": table_headers, "rows": table_rows})

        # Extract images from inline shapes and relationships
        try:
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        img_data = rel.target_part.blob
                        if img_data and len(img_data) <= 2 * 1024 * 1024:
                            ext = rel.target_ref.split(".")[-1] if "." in rel.target_ref else "png"
                            all_images.append({
                                "data_b64": base64.b64encode(img_data).decode("ascii"),
                                "page": 1,
                                "width": 0,
                                "height": 0,
                                "ext": ext,
                            })
                            md_parts.append(f"![Image {len(all_images)}](embedded)")
                    except Exception:
                        continue
        except Exception:
            pass

        # Extract footnotes and endnotes
        footnotes_md = _extract_docx_notes(raw_bytes)
        if footnotes_md:
            md_parts.append("\n## Notes\n")
            md_parts.append(footnotes_md)

        full_text = "\n".join(text_parts)
        word_count = len(full_text.split())
        markdown = "\n\n".join(md_parts)

        metadata["word_count"] = word_count
        metadata["paragraph_count"] = len(doc.paragraphs)
        metadata["table_count"] = len(doc.tables)
        metadata["image_count"] = len(all_images)

        return DocumentResult(
            text=full_text,
            markdown=markdown,
            metadata=metadata,
            page_count=1,
            word_count=word_count,
            tables=all_tables,
            images=all_images,
        )

    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return DocumentResult(
            text=f"[DOCX extraction failed: {e}]",
            markdown=f"*DOCX extraction failed: {e}*",
            metadata={"error": str(e), "document_type": "docx"},
        )


def _build_formatted_runs(para, hyperlink_map: dict) -> str:
    """Build markdown-formatted text from paragraph runs with bold/italic/hyperlinks."""
    from lxml import etree
    nsmap = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
             "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"}

    parts = []
    for child in para._element:
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ""

        if tag == "r":
            # Normal run
            run_text = ""
            is_bold = False
            is_italic = False
            for sub in child:
                sub_tag = etree.QName(sub.tag).localname if isinstance(sub.tag, str) else ""
                if sub_tag == "rPr":
                    for prop in sub:
                        prop_tag = etree.QName(prop.tag).localname if isinstance(prop.tag, str) else ""
                        if prop_tag == "b":
                            val = prop.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val")
                            is_bold = val != "0" if val else True
                        elif prop_tag == "i":
                            val = prop.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val")
                            is_italic = val != "0" if val else True
                elif sub_tag == "t":
                    run_text += sub.text or ""

            if run_text:
                if is_bold and is_italic:
                    parts.append(f"***{run_text}***")
                elif is_bold:
                    parts.append(f"**{run_text}**")
                elif is_italic:
                    parts.append(f"*{run_text}*")
                else:
                    parts.append(run_text)

        elif tag == "hyperlink":
            # Hyperlink element
            r_id = child.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id", "")
            link_url = hyperlink_map.get(r_id, "")
            link_text = ""
            for sub in child.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                link_text += sub.text or ""
            if link_text and link_url:
                parts.append(f"[{link_text}]({link_url})")
            elif link_text:
                parts.append(link_text)

    return "".join(parts) if parts else para.text.strip()


def _extract_docx_notes(raw_bytes: bytes) -> str:
    """Extract footnotes and endnotes from DOCX ZIP package."""
    notes = []
    try:
        with zipfile.ZipFile(io.BytesIO(raw_bytes)) as zf:
            for note_file in ("word/footnotes.xml", "word/endnotes.xml"):
                if note_file not in zf.namelist():
                    continue
                from lxml import etree
                tree = etree.parse(zf.open(note_file))
                ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                note_tag = "footnote" if "footnote" in note_file else "endnote"
                for note in tree.findall(f".//w:{note_tag}", ns):
                    note_id = note.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}id", "")
                    # Skip separator notes (id 0 and -1)
                    if note_id in ("0", "-1"):
                        continue
                    text_parts = []
                    for t in note.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"):
                        if t.text:
                            text_parts.append(t.text)
                    text = "".join(text_parts).strip()
                    if text:
                        notes.append(f"[^{note_id}]: {text}")
    except Exception:
        pass
    return "\n".join(notes)


# ---------------------------------------------------------------------------
# XLSX Extraction
# ---------------------------------------------------------------------------

async def extract_xlsx(raw_bytes: bytes) -> DocumentResult:
    """Extract data from XLSX using openpyxl."""
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl not installed, returning empty result")
        return DocumentResult(
            text="[XLSX extraction requires openpyxl]",
            markdown="*XLSX extraction requires openpyxl*",
            metadata={"error": "openpyxl not installed"},
        )

    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)

        metadata = {
            "document_type": "xlsx",
            "sheet_names": wb.sheetnames,
            "sheet_count": len(wb.sheetnames),
        }

        md_parts = []
        text_parts = []
        all_tables = []
        total_rows = 0
        has_formulas = False

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            md_parts.append(f"## Sheet: {sheet_name}\n")

            rows_data = []
            for row in ws.iter_rows(values_only=True):
                row_vals = []
                for cell in row:
                    if cell is None:
                        row_vals.append("")
                    else:
                        val = str(cell)
                        if val.startswith("="):
                            has_formulas = True
                        row_vals.append(val)
                rows_data.append(row_vals)
                total_rows += 1

            if not rows_data:
                md_parts.append("*Empty sheet*\n")
                continue

            # First row as headers
            headers = rows_data[0]
            data_rows = rows_data[1:]

            all_tables.append({
                "headers": headers,
                "rows": data_rows,
                "sheet": sheet_name,
            })

            md_parts.append(_table_to_markdown(headers, data_rows))
            md_parts.append("")

            for row in rows_data:
                text_parts.append("\t".join(row))

        wb.close()

        full_text = "\n".join(text_parts)
        word_count = len(full_text.split())
        markdown = "\n\n".join(md_parts)

        metadata["total_rows"] = total_rows
        metadata["has_formulas"] = has_formulas

        return DocumentResult(
            text=full_text,
            markdown=markdown,
            metadata=metadata,
            page_count=len(wb.sheetnames) if hasattr(wb, 'sheetnames') else metadata["sheet_count"],
            word_count=word_count,
            tables=all_tables,
        )

    except Exception as e:
        logger.error(f"XLSX extraction failed: {e}")
        return DocumentResult(
            text=f"[XLSX extraction failed: {e}]",
            markdown=f"*XLSX extraction failed: {e}*",
            metadata={"error": str(e), "document_type": "xlsx"},
        )


# ---------------------------------------------------------------------------
# PPTX Extraction
# ---------------------------------------------------------------------------

async def extract_pptx(raw_bytes: bytes) -> DocumentResult:
    """Extract text, tables, and notes from PPTX using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed, returning empty result")
        return DocumentResult(
            text="[PPTX extraction requires python-pptx]",
            markdown="*PPTX extraction requires python-pptx*",
            metadata={"error": "python-pptx not installed"},
        )

    try:
        prs = Presentation(io.BytesIO(raw_bytes))

        metadata = {
            "document_type": "pptx",
            "slide_count": len(prs.slides),
        }

        md_parts = []
        text_parts = []
        all_tables = []
        all_images = []
        has_notes = False

        for slide_num, slide in enumerate(prs.slides, 1):
            # Get slide title
            slide_title = ""
            if slide.shapes.title:
                slide_title = slide.shapes.title.text.strip()

            md_parts.append(f"## Slide {slide_num}" + (f": {slide_title}" if slide_title else ""))
            md_parts.append("")

            for shape in slide.shapes:
                # Text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)
                            md_parts.append(text)

                # Tables
                if shape.has_table:
                    table = shape.table
                    headers = [cell.text.strip() for cell in table.rows[0].cells]
                    rows = []
                    for row_idx in range(1, len(table.rows)):
                        row_cells = [cell.text.strip() for cell in table.rows[row_idx].cells]
                        rows.append(row_cells)
                    all_tables.append({"headers": headers, "rows": rows, "slide": slide_num})
                    md_parts.append("")
                    md_parts.append(_table_to_markdown(headers, rows))
                    md_parts.append("")

                # Images
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        img_blob = shape.image.blob
                        if img_blob and len(img_blob) <= 2 * 1024 * 1024:
                            ext = shape.image.content_type.split("/")[-1] if shape.image.content_type else "png"
                            all_images.append({
                                "data_b64": base64.b64encode(img_blob).decode("ascii"),
                                "page": slide_num,
                                "width": shape.width,
                                "height": shape.height,
                                "ext": ext,
                            })
                            md_parts.append(f"![Image from slide {slide_num}](embedded)")
                    except Exception:
                        pass

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    has_notes = True
                    md_parts.append(f"\n> Note: {notes_text}")
                    text_parts.append(f"[Speaker Note] {notes_text}")

            md_parts.append("")

        metadata["has_notes"] = has_notes
        metadata["table_count"] = len(all_tables)
        metadata["image_count"] = len(all_images)

        full_text = "\n".join(text_parts)
        word_count = len(full_text.split())
        markdown = "\n\n".join(md_parts)

        return DocumentResult(
            text=full_text,
            markdown=markdown,
            metadata=metadata,
            page_count=len(prs.slides),
            word_count=word_count,
            tables=all_tables,
            images=all_images,
        )

    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return DocumentResult(
            text=f"[PPTX extraction failed: {e}]",
            markdown=f"*PPTX extraction failed: {e}*",
            metadata={"error": str(e), "document_type": "pptx"},
        )


# ---------------------------------------------------------------------------
# CSV Extraction
# ---------------------------------------------------------------------------

async def extract_csv(raw_bytes: bytes) -> DocumentResult:
    """Extract data from CSV/TSV files with auto-detection of delimiter and encoding."""
    # Try encoding detection
    text_content = None
    encoding_used = "utf-8"
    for encoding in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            text_content = raw_bytes.decode(encoding)
            encoding_used = encoding
            break
        except (UnicodeDecodeError, ValueError):
            continue

    if text_content is None:
        return DocumentResult(
            text="[CSV decoding failed]",
            markdown="*CSV decoding failed — unsupported encoding*",
            metadata={"error": "encoding detection failed", "document_type": "csv"},
        )

    try:
        # Auto-detect delimiter
        delimiter = ","
        try:
            sample = text_content[:8192]
            dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
            delimiter = dialect.delimiter
        except csv.Error:
            pass

        reader = csv.reader(io.StringIO(text_content), delimiter=delimiter)
        all_rows = list(reader)

        if not all_rows:
            return DocumentResult(
                text="[Empty CSV]",
                markdown="*Empty CSV file*",
                metadata={"document_type": "csv", "row_count": 0},
            )

        headers = all_rows[0]
        data_rows = all_rows[1:]
        total_rows = len(data_rows)

        metadata = {
            "document_type": "csv",
            "row_count": total_rows,
            "column_count": len(headers),
            "delimiter": repr(delimiter),
            "encoding": encoding_used,
        }

        # For large CSVs, truncate display
        if total_rows > 1000:
            display_rows = data_rows[:500] + [["..."] * len(headers)] + data_rows[-100:]
            metadata["truncated"] = True
            metadata["display_note"] = f"Showing first 500 + last 100 of {total_rows} rows"
        else:
            display_rows = data_rows

        md_parts = [f"## CSV Data ({total_rows} rows, {len(headers)} columns)\n"]
        md_parts.append(_table_to_markdown(headers, display_rows))

        text_parts = [delimiter.join(headers)]
        for row in data_rows:
            text_parts.append(delimiter.join(row))

        full_text = "\n".join(text_parts)
        word_count = len(full_text.split())
        markdown = "\n\n".join(md_parts)

        return DocumentResult(
            text=full_text,
            markdown=markdown,
            metadata=metadata,
            page_count=1,
            word_count=word_count,
            tables=[{"headers": headers, "rows": data_rows}],
        )

    except Exception as e:
        logger.error(f"CSV extraction failed: {e}")
        return DocumentResult(
            text=f"[CSV extraction failed: {e}]",
            markdown=f"*CSV extraction failed: {e}*",
            metadata={"error": str(e), "document_type": "csv"},
        )


# ---------------------------------------------------------------------------
# RTF Extraction
# ---------------------------------------------------------------------------

async def extract_rtf(raw_bytes: bytes) -> DocumentResult:
    """Extract text from RTF using striprtf."""
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        logger.warning("striprtf not installed, returning empty result")
        return DocumentResult(
            text="[RTF extraction requires striprtf]",
            markdown="*RTF extraction requires striprtf*",
            metadata={"error": "striprtf not installed"},
        )

    try:
        # Decode RTF content
        rtf_content = None
        for encoding in ("utf-8", "latin-1", "cp1252"):
            try:
                rtf_content = raw_bytes.decode(encoding)
                break
            except (UnicodeDecodeError, ValueError):
                continue

        if rtf_content is None:
            return DocumentResult(
                text="[RTF decoding failed]",
                metadata={"error": "encoding detection failed", "document_type": "rtf"},
            )

        text = rtf_to_text(rtf_content)
        word_count = len(text.split())

        # Basic markdown conversion — split into paragraphs
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        markdown = "\n\n".join(paragraphs)

        return DocumentResult(
            text=text,
            markdown=markdown,
            metadata={"document_type": "rtf", "word_count": word_count},
            page_count=1,
            word_count=word_count,
        )

    except Exception as e:
        logger.error(f"RTF extraction failed: {e}")
        return DocumentResult(
            text=f"[RTF extraction failed: {e}]",
            markdown=f"*RTF extraction failed: {e}*",
            metadata={"error": str(e), "document_type": "rtf"},
        )


# ---------------------------------------------------------------------------
# EPUB Extraction
# ---------------------------------------------------------------------------

async def extract_epub(raw_bytes: bytes) -> DocumentResult:
    """Extract text from EPUB using zipfile + BeautifulSoup (no extra dependency)."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return DocumentResult(
            text="[EPUB extraction requires beautifulsoup4]",
            metadata={"error": "beautifulsoup4 not installed"},
        )

    try:
        with zipfile.ZipFile(io.BytesIO(raw_bytes)) as zf:
            # Parse container.xml to find OPF
            container = zf.read("META-INF/container.xml")
            from lxml import etree
            container_tree = etree.fromstring(container)
            ns = {"c": "urn:oasis:names:tc:opendocument:xmlns:container"}
            rootfile = container_tree.find(".//c:rootfile", ns)
            if rootfile is None:
                return DocumentResult(
                    text="[Invalid EPUB: no rootfile]",
                    metadata={"error": "no rootfile in container.xml", "document_type": "epub"},
                )
            opf_path = rootfile.get("full-path", "")

            # Parse OPF for metadata and spine
            opf_content = zf.read(opf_path)
            opf_tree = etree.fromstring(opf_content)
            opf_ns = {
                "opf": "http://www.idpf.org/2007/opf",
                "dc": "http://purl.org/dc/elements/1.1/",
            }

            # Extract metadata
            metadata = {"document_type": "epub"}
            title_el = opf_tree.find(".//dc:title", opf_ns)
            if title_el is not None and title_el.text:
                metadata["title"] = title_el.text
            author_el = opf_tree.find(".//dc:creator", opf_ns)
            if author_el is not None and author_el.text:
                metadata["author"] = author_el.text
            lang_el = opf_tree.find(".//dc:language", opf_ns)
            if lang_el is not None and lang_el.text:
                metadata["language"] = lang_el.text
            publisher_el = opf_tree.find(".//dc:publisher", opf_ns)
            if publisher_el is not None and publisher_el.text:
                metadata["publisher"] = publisher_el.text

            # Get spine order → manifest id → href
            manifest_items = {}
            for item in opf_tree.findall(".//opf:manifest/opf:item", opf_ns):
                item_id = item.get("id", "")
                href = item.get("href", "")
                media_type = item.get("media-type", "")
                manifest_items[item_id] = {"href": href, "media_type": media_type}

            spine_ids = []
            for itemref in opf_tree.findall(".//opf:spine/opf:itemref", opf_ns):
                idref = itemref.get("idref", "")
                if idref:
                    spine_ids.append(idref)

            # Resolve base directory for OPF-relative paths
            import posixpath
            opf_dir = posixpath.dirname(opf_path)

            md_parts = []
            text_parts = []
            all_images = []
            chapter_count = 0

            if metadata.get("title"):
                md_parts.append(f"# {metadata['title']}\n")
            if metadata.get("author"):
                md_parts.append(f"**Author:** {metadata['author']}\n")
            md_parts.append("---\n")

            # Extract cover image if present
            for item_id, item_info in manifest_items.items():
                if "cover" in item_id.lower() and "image" in item_info["media_type"]:
                    try:
                        cover_path = posixpath.join(opf_dir, item_info["href"]) if opf_dir else item_info["href"]
                        cover_data = zf.read(cover_path)
                        if len(cover_data) <= 2 * 1024 * 1024:
                            ext = item_info["href"].split(".")[-1]
                            all_images.append({
                                "data_b64": base64.b64encode(cover_data).decode("ascii"),
                                "page": 0,
                                "width": 0,
                                "height": 0,
                                "ext": ext,
                            })
                    except Exception:
                        pass
                    break

            # Read content files in spine order
            for spine_id in spine_ids:
                item = manifest_items.get(spine_id)
                if not item or "html" not in item["media_type"].lower():
                    continue

                href = posixpath.join(opf_dir, item["href"]) if opf_dir else item["href"]
                try:
                    html_content = zf.read(href).decode("utf-8", errors="replace")
                except (KeyError, Exception):
                    continue

                soup = BeautifulSoup(html_content, "lxml")

                # Remove scripts and styles
                for tag in soup(["script", "style"]):
                    tag.decompose()

                chapter_text = soup.get_text(separator="\n", strip=True)
                if not chapter_text.strip():
                    continue

                chapter_count += 1
                text_parts.append(chapter_text)

                # Convert HTML to markdown using markdownify if available
                try:
                    from markdownify import markdownify as md
                    chapter_md = md(str(soup.body or soup), heading_style="ATX", strip=["script", "style"])
                except ImportError:
                    chapter_md = chapter_text

                md_parts.append(f"## Chapter {chapter_count}\n")
                md_parts.append(chapter_md.strip())
                md_parts.append("\n---\n")

            metadata["chapter_count"] = chapter_count
            metadata["image_count"] = len(all_images)

            full_text = "\n\n".join(text_parts)
            word_count = len(full_text.split())
            markdown = "\n\n".join(md_parts)

            return DocumentResult(
                text=full_text,
                markdown=markdown,
                metadata=metadata,
                page_count=chapter_count,
                word_count=word_count,
                images=all_images,
            )

    except Exception as e:
        logger.error(f"EPUB extraction failed: {e}")
        return DocumentResult(
            text=f"[EPUB extraction failed: {e}]",
            markdown=f"*EPUB extraction failed: {e}*",
            metadata={"error": str(e), "document_type": "epub"},
        )


# ---------------------------------------------------------------------------
# Unified dispatch
# ---------------------------------------------------------------------------

async def extract_document(raw_bytes: bytes, doc_type: str) -> DocumentResult:
    """Route to the correct extractor based on document type."""
    extractors = {
        "pdf": extract_pdf,
        "docx": extract_docx,
        "xlsx": extract_xlsx,
        "pptx": extract_pptx,
        "csv": extract_csv,
        "rtf": extract_rtf,
        "epub": extract_epub,
    }
    extractor = extractors.get(doc_type)
    if not extractor:
        return DocumentResult(
            text=f"[Unsupported document type: {doc_type}]",
            metadata={"error": f"Unsupported: {doc_type}"},
        )
    return await extractor(raw_bytes)
